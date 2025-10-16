import sys
import json
import re
from pptx import Presentation

def is_caption_like(text):
    if not text or len(text.strip()) > 100:
        return False

    caption_patterns = [
        r'^Figure\s+\d+(\.\d+)?[:\-]?',
        r'^Chart\s+\d+[:\-]?',
        r'^Image\s+\d+[:\-]?',
        r'^Diagram\s+\d+[:\-]?',
        r'^\(?Figure\s+\d+.*\)?$',
        r'^Fig\.?\s+\d+.*',
    ]
    return any(re.match(pat, text.strip()) for pat in caption_patterns)

def extract_text_from_pptx(file_path):
    try:
        prs = Presentation(file_path)
        markdown = ""

        for slide in prs.slides:
            slide_text = []

            #note: add this prompt for revieweing captions before ignoring --> "Ignore lines starting with [Caption] unless they include relevant definitions or explanations."
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if not text:
                        continue
                    if is_caption_like(text):
                        slide_text.append(f"[Caption] {text}")
                    else:
                        slide_text.append(text)

            if slide_text:
                markdown += "\n\n" + "\n".join(slide_text)

        markdown = markdown.strip()

        result = {
            "success": True,
            "markdown": markdown
        }

        print(json.dumps(result))

    except Exception as e:
        print(json.dumps({ "success": False, "error": str(e) }))

if __name__ == "__main__":
    file_path = sys.argv[1]
    extract_text_from_pptx(file_path)
